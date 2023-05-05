import openai
import json5
import pprint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import os
import requests
import random
import string

class GeneratePPT:
    def __init__(self, openai_key, topic, slide_count, output_filename, access_key):
        self.openai_key  = openai_key
        self.topic       = str(topic) 
        self.slide_count = int(slide_count) 
        self.output_filename = output_filename
        self.access_key = access_key

    def start(self):
        content_gen = ContentGenerator(self.topic, self.slide_count, self.openai_key)
        json_data = content_gen.start()
        
        ppt = PPTGenerator(json_data, self.output_filename, self.access_key)
        ppt.generate_ppt()

class ContentGenerator:
    def __init__(self, topic, slide_count, openai_key):
        # Set up OpenAI API credentials
        openai.api_key = openai_key

        self.topic       = str(topic)
        self.slide_count = str(slide_count)
        self.prompt = "write content for ppt, topic is 'TOPIC' and slide count is SLIDECOUNT, give me the heading and content of each slide, give the output in json format, where there should be a key with the name 'Content' which is equal to list of dictionary, dictionary should contain two keys, one is 'Heading' & 'SlideContent', slide content, should be the list of points (string)"

    def start(self):
        self.current_prompt = self.prompt.replace("TOPIC", self.topic).replace("SLIDECOUNT", self.slide_count)
        json_data =self.generate_powerpoint_content()
        pprint.pprint(json_data)
        return json_data

    def generate_powerpoint_content(self):
        response = openai.Completion.create(
            engine="text-davinci-002",
            prompt=self.current_prompt,
            max_tokens=2048,
            n=1,
            stop=None,
            temperature=0.5,
        )

        output_text = response.choices[0].text
        json_obj = json5.loads(output_text)
        """
        {Content': [
            {'Heading': 'What is Waste Management?',
             'SlideContent': 'Waste management is the process of treating '}]
        }
        """
        return json_obj

class PPTGenerator:
    def __init__(self, json_data, ppt_file, access_key):
        self.json_data = json_data
        self.ppt_file = ppt_file
        self.access_key = access_key
        self.prs = Presentation()

    def generate_ppt(self):  
        index = 0
        image_getter = ImageGetter(self.access_key)
        for data in self.json_data["Content"]:
            heading = data["Heading"]
            content = data["SlideContent"]
            image_file = image_getter.get_image_url(heading)
            self.add_slide(image_file, heading, content)
            index += 1
            os.remove(image_file)   

        # Save the PowerPoint presentation
        self.prs.save(self.ppt_file) 

    def add_slide(self, image_file, heading, content):
        slide_layout = self.prs.slide_layouts[6]  # Blank layout

        # Add slide with Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Set image as slide background
        slide.shapes.add_picture(image_file, 0, 0, self.prs.slide_width, self.prs.slide_height)

        # Add heading to slide
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.5)
        shape = slide.shapes.add_textbox(left, top, width, height)
        para = shape.text_frame.add_paragraph()
        para.text = heading
        para.font.bold = True
        para.font.underline = True
        para.alignment = MSO_ANCHOR.MIDDLE  # Horizontally center the text
        para.vertical_anchor = MSO_ANCHOR.TOP  # Align the text to the top of the text box
        para.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Add content to slide
        left = Inches(0.5)
        width = Inches(9)
        height = Inches(2)

        top_height = 1.5
        for paragraph in content:
            top = Inches(top_height)
            shape = slide.shapes.add_textbox(left, top, width, height)
            text_frame = shape.text_frame
            lines_list  = []
            word_cout = 0
            line = ""
            for word in paragraph.split():
                word_cout += 1
                line += word + " "
                if word_cout == 6:
                    lines_list.append(line)
                    word_cout = 0
                    line = ""

            for line in lines_list:
                # Add a line of text
                p = text_frame.add_paragraph()
                p.level = 0
                p.text = line
                p.font.size = Pt(16)
                p.font.bullet = True
                p.font.bullet_char = '•'
                p.font.color.rgb = RGBColor(255,255,255)
                top_height += 0.4

class ImageGetter:
    def __init__(self, access_key):
        self.access_key = access_key
        self.base_url = "https://api.unsplash.com/photos/random"
        self.params = {
            "orientation": "landscape",
            "query": "",
            "w": 256, #1920, 
            "h": 144, #1080,
            "client_id": self.access_key
        }

    def get_image_url(self, topic):
        self.params["query"] = topic
        response = requests.get(self.base_url, params=self.params)
        letters = string.ascii_lowercase
        filename = ''.join(random.choice(letters) for i in range(6))
        if response.status_code == 200:
            return self.download(response.json()["urls"]["raw"], "static/"+filename+".jpg")
        else:
            return None
        
    def download(self, url, filename):
        r = requests.get(url, allow_redirects=True)
        open(filename, 'wb').write(r.content)
        return filename


if __name__ == "__main__":
    topic = "Health"
    slide_count = 2
    output_filename = "output.pptx"
    openai_key  = "sk-u7yHXXXXXXXXXXXXXXXXXXXXXXXXXXXen5ji0hpTuBs8Z9"
    access_key = "YRx46BWXXXXXXXXXXXXXXXXXXXXXXXXXXNA6jecP0-91Syw"
    test = GeneratePPT(openai_key, topic, slide_count, output_filename, access_key)
    test.start()

